"""
Build report/slides.pptx -- presentation deck for the 10-12 minute talk.

Design principles applied here:
  * one main idea per slide (50 s of spoken content, not a wall of text)
  * formulas kept where they earn their place, embedded as crisp PNGs
    rendered by report/build_docx.py (so they look like LaTeX, not
    Unicode strings)
  * visual hierarchy: coloured title band, large bullets, occasional
    rounded-rectangle "cards" for side-by-side concepts
  * font sizes designed to read from the back of a classroom
  * plot images embedded full-bleed where the figure IS the content

Run from the repo root:
    python report/build_pptx.py
"""

from __future__ import annotations

import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG_DIR   = os.path.join(REPO_ROOT, "figures")
EQ_DIR    = os.path.join(REPO_ROOT, "report", "eqs")
OUT_PATH  = os.path.join(REPO_ROOT, "report", "slides.pptx")


# ---- Theme colours ----
NAVY      = RGBColor(0x1F, 0x4E, 0x79)
RED       = RGBColor(0xC0, 0x50, 0x4D)
GREEN     = RGBColor(0x2C, 0xA0, 0x2C)
GREY      = RGBColor(0x7F, 0x7F, 0x7F)
DARK      = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT     = RGBColor(0xF4, 0xF4, 0xF6)
ACCENT    = RGBColor(0xE5, 0xC1, 0x07)
SOFT_NAVY = RGBColor(0xE5, 0xEC, 0xF3)
SOFT_RED  = RGBColor(0xFA, 0xE7, 0xE6)
SOFT_GRN  = RGBColor(0xE0, 0xF1, 0xE0)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_textbox(slide, left, top, width, height, text, *,
                size=20, bold=False, italic=False, color=DARK,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                size=22, color=DARK, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, txt in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(12)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = "•  " + txt
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def add_title_band(slide, title, subtitle=None):
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.95))
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.shadow.inherit = False

    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.10), Inches(12.3), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"
    if subtitle:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(13)
        run2.font.italic = True
        run2.font.color.rgb = RGBColor(0xE0, 0xE5, 0xEC)
        run2.font.name = "Calibri"


def add_card(slide, left, top, width, height, *,
             fill=SOFT_NAVY, border=NAVY):
    """Coloured rounded-rectangle card for side-by-side concepts."""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(1.25)
    card.shadow.inherit = False
    return card


def add_card_text(slide, left, top, width, height, header, body, *,
                  header_color=NAVY, body_color=DARK,
                  fill=SOFT_NAVY, border=NAVY):
    add_card(slide, left, top, width, height, fill=fill, border=border)
    add_textbox(slide, left + 0.2, top + 0.2, width - 0.4, 0.5,
                header, size=20, bold=True, color=header_color)
    add_textbox(slide, left + 0.2, top + 0.85, width - 0.4, height - 1.0,
                body, size=16, color=body_color)


def add_figure(slide, fname, left, top, width=None, height=None):
    path = os.path.join(FIG_DIR, fname)
    if not os.path.isfile(path):
        print(f"[WARN] missing figure: {path}")
        return
    kw = {}
    if width  is not None: kw["width"]  = Inches(width)
    if height is not None: kw["height"] = Inches(height)
    slide.shapes.add_picture(path, Inches(left), Inches(top), **kw)


def add_eq_image(slide, name, left, top, width=None, height=None):
    path = os.path.join(EQ_DIR, f"{name}.png")
    if not os.path.isfile(path):
        print(f"[WARN] missing equation: {path}")
        return
    kw = {}
    if width  is not None: kw["width"]  = Inches(width)
    if height is not None: kw["height"] = Inches(height)
    slide.shapes.add_picture(path, Inches(left), Inches(top), **kw)


def add_footer(slide, page_num, total=13):
    add_textbox(slide, 11.0, 7.10, 2.2, 0.3,
                f"Slide {page_num} / {total}",
                size=10, color=GREY, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Build the deck
# ---------------------------------------------------------------------------

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

TOTAL = 13


# ---- Slide 1: Title ---------------------------------------------------
s = add_blank(prs)
# soft accent strip on the left
strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                           Inches(0), Inches(0), Inches(0.4), Inches(7.5))
strip.line.fill.background(); strip.fill.solid()
strip.fill.fore_color.rgb = ACCENT

add_textbox(s, 1.2, 1.7, 11.0, 0.5,
            "CY315  ·  Wireless and Mobile Security",
            size=14, italic=True, color=GREY)

add_textbox(s, 1.2, 2.3, 11.0, 1.4,
            "Defending the Wiretap Channel",
            size=44, bold=True, color=NAVY)

add_textbox(s, 1.2, 3.6, 11.0, 1.0,
            "A deep-RL approach to artificial-noise power allocation"
            "  —  when Alice's intel on Eve is noisy",
            size=22, italic=True, color=DARK)

add_textbox(s, 1.2, 5.4, 11.0, 0.5,
            "Muhammad Ismail   ·   Abubakar Butt   ·   Usman Ali",
            size=18, color=DARK)
add_textbox(s, 1.2, 5.95, 11.0, 0.4,
            "2023453   ·   2023352   ·   2023581",
            size=13, italic=True, color=GREY)


# ---- Slide 2: The problem (story-first, no formulas) -----------------
s = add_blank(prs)
add_title_band(s, "The problem we're solving")

add_textbox(s, 0.6, 1.2, 12.3, 0.6,
            "Wireless is broadcast.  Anything Alice transmits, "
            "Eve picks up too.",
            size=22, italic=True, color=NAVY)

add_bullets(s, 0.6, 2.0, 12.3, 5.0, [
    "Standard fix: artificial noise (AN) sprayed into Bob's null "
    "space  —  Bob hears clean signal, Eve hears signal + jam.",
    "The knob: how much power to put on the message vs the jamming.",
    "The classic answer (Goel & Negi, 2008): split it 50/50.  Safe but "
    "leaves performance on the table.",
    "The textbook upgrade: optimise the split per channel.  Works "
    "great  —  IF Alice has perfect information about Eve.  She "
    "usually doesn't.",
    "Plug noisy intel into the textbook optimiser  →  a confidently-"
    "wrong answer that performs WORSE than doing nothing.  We'll see "
    "this on the next plot.",
], size=20)
add_footer(s, 2)


# ---- Slide 3: System model (geometry sketch hero + clean math row) --
s = add_blank(prs)
add_title_band(s, "The setting  —  MISO wiretap channel",
               "Alice (multi-antenna), Bob (legitimate), Eve (eavesdropper)")

# Hero: the geometry sketch, centred, with breathing room.
add_figure(s, "14_geometry_sketch.png",
           left=1.6, top=1.05, width=10.1)

# Single horizontal row at the bottom: math card on the left,
# kappa legend card on the right. Plenty of vertical and horizontal
# spacing so nothing crowds anything else.
add_card(s, 0.5, 5.65, 6.4, 1.65, fill=SOFT_NAVY, border=NAVY)
add_textbox(s, 0.65, 5.75, 6.1, 0.4,
            "Imperfect CSI model",
            size=13, bold=True, color=NAVY)
add_eq_image(s, "eq04_csi", left=0.85, top=6.20, width=5.7)

add_card(s, 7.1, 5.65, 5.7, 1.65, fill=SOFT_GRN, border=GREEN)
add_textbox(s, 7.25, 5.75, 5.4, 0.4,
            "κ  —  CSI quality dial",
            size=13, bold=True, color=GREEN)
add_textbox(s, 7.25, 6.20, 5.4, 1.0,
            "κ = 1  perfect intel    ·    "
            "κ = 0.4  our default    ·    "
            "κ = 0  pure noise",
            size=12, color=DARK)

add_footer(s, 3)


# ---- Slide 4: The two classical answers ------------------------------
s = add_blank(prs)
add_title_band(s, "Two classical answers  —  and why both fall short")

# Card 1: Fixed
add_card_text(s, 0.6, 1.5, 6.0, 5.5,
    "Scheme 1  —  Fixed (Goel & Negi, 2008)",
    "Always splits power 50/50.\n\n"
    "Pros\n"
    "    Cannot be misled by bad intel.\n"
    "    Always returns a sensible answer.\n\n"
    "Cons\n"
    "    Sub-optimal whenever the channel is asymmetric.\n"
    "    Leaves rate on the table at high SNR.",
    fill=SOFT_NAVY, border=NAVY, header_color=NAVY)

# Card 2: Traditional
add_card_text(s, 6.7, 1.5, 6.1, 5.5,
    "Scheme 2  —  Traditional optimiser",
    "Per-channel argmax of Rs against the noisy ĥE.\n"
    "scipy bounded scalar (Brent's method).\n\n"
    "Pros\n"
    "    Provably optimal at κ = 1.\n\n"
    "Cons\n"
    "    Trusts ĥE blindly.  At low κ the noise gets amplified, "
    "and the optimiser confidently picks the WRONG ρ.",
    fill=SOFT_RED, border=RED, header_color=RED)

add_footer(s, 4)


# ---- Slide 5: BEFORE vs AFTER our AI (full-bleed comparison) ---------
s = add_blank(prs)
add_title_band(s, "BEFORE vs AFTER our AI  —  same channels, same κ",
               "Left: only existing schemes.  Right: with DQN added.")

# Hero figure: the new before/after two-panel comparison
add_figure(s, "12_before_after.png",
           left=0.3, top=1.15, width=12.7)

# Axis-callout strip below the figure — large, explicit
ax_band = s.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.55))
ax_band.line.color.rgb = NAVY
ax_band.line.width = Pt(0.75)
ax_band.fill.solid()
ax_band.fill.fore_color.rgb = SOFT_NAVY
ax_band.shadow.inherit = False

add_textbox(s, 0.7, 6.55, 6.0, 0.45,
            "X axis  →  Transmit power, SNR in dB  (0 to 30)",
            size=13, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
add_textbox(s, 6.5, 6.55, 6.3, 0.45,
            "Y axis  →  Average secrecy rate  Rs (bits/s/Hz)  —  "
            "higher is better",
            size=13, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

# Story strip under the axis row
add_textbox(s, 0.5, 7.10, 12.3, 0.35,
            "Left panel: Traditional (red) drops BELOW Fixed (grey) — "
            "trusting noisy intel hurts.    "
            "Right panel: DQN (green) holds at or above Fixed everywhere.",
            size=12, italic=True, color=DARK, align=PP_ALIGN.CENTER)

add_footer(s, 5)


# ---- Slide 6: Our approach (concept, not architecture) ---------------
s = add_blank(prs)
add_title_band(s, "Our approach  —  train a policy that EXPECTS noisy intel")

# Three idea-blocks
add_card(s, 0.6, 1.4, 4.0, 5.5, fill=SOFT_NAVY, border=NAVY)
add_textbox(s, 0.6, 1.6, 4.0, 0.5, "1.  Reframe",
            size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_textbox(s, 0.8, 2.2, 3.6, 4.5,
            "Stop trying to find the closed-form optimum each call.\n\n"
            "Instead, train a policy off-line, across MANY noisy-intel "
            "scenarios, that learns when to trust ĥE and when to "
            "fall back to safe defaults.",
            size=15, color=DARK)

add_card(s, 4.7, 1.4, 4.0, 5.5, fill=SOFT_GRN, border=GREEN)
add_textbox(s, 4.7, 1.6, 4.0, 0.5, "2.  Two key inputs",
            size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_textbox(s, 4.9, 2.2, 3.6, 4.5,
            "Tell the agent how trustworthy its intel is  →  feed κ "
            "directly into its observation.\n\n"
            "Tell it the channel geometry  →  feed the alignment "
            "between Bob's beam and the noisy Eve estimate.",
            size=15, color=DARK)

add_card(s, 8.8, 1.4, 4.0, 5.5, fill=SOFT_RED, border=RED)
add_textbox(s, 8.8, 1.6, 4.0, 0.5, "3.  Single deployable model",
            size=22, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_textbox(s, 9.0, 2.2, 3.6, 4.5,
            "ONE trained network handles every (SNR, κ) it sees in the "
            "test set.\n\n"
            "No retuning, no re-optimising at run time.  A forward pass "
            "replaces an iterative solver.",
            size=15, color=DARK)

add_footer(s, 6)


# ---- Slide 7: What the agent sees (icons + ONE equation) -------------
s = add_blank(prs)
add_title_band(s, "What the agent sees, what it does")

# Left: state items as icons + labels (no full formula)
add_textbox(s, 0.6, 1.2, 6.5, 0.5,
            "State  —  7 numbers describing the situation",
            size=20, bold=True, color=NAVY)

state_items = [
    ("📡", "Bob's channel strength"),
    ("📡", "Eve's (noisy) channel strength"),
    ("⚡",  "Transmit SNR"),
    ("🕒", "What it picked last time"),
    ("📊", "How it scored last time"),
    ("🎯", "κ  —  intel quality"),
    ("📐", "α  —  beam-vs-Eve alignment"),
]
y = 1.85
for icon, label in state_items:
    add_textbox(s, 0.7, y, 0.4, 0.35, icon, size=16,
                color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(s, 1.2, y, 5.5, 0.35, label, size=15, color=DARK)
    y += 0.42

# small alignment-formula image as a sidebar
add_textbox(s, 0.7, 5.6, 6.0, 0.35,
            "(α formally:)",
            size=12, italic=True, color=GREY)
add_eq_image(s, "eq06_alpha", left=0.7, top=5.95, width=5.5)

# Right: action + network (text-card style)
add_card_text(s, 7.5, 1.2, 5.3, 2.5,
    "Action",
    "Pick one of 17 candidate splits between message and AN  "
    "(steps of 0.05).\n\n"
    "Discrete  →  light DQN architecture, fast inference.",
    fill=SOFT_NAVY, border=NAVY, header_color=NAVY)

add_card_text(s, 7.5, 3.9, 5.3, 3.0,
    "Network",
    "Small feed-forward net, 7 inputs  →  17 Q-values.\n"
    "Greedy at test time:  pick the highest-Q action.\n\n"
    "Trained with Huber loss + Adam over 7000 episodes.",
    fill=SOFT_GRN, border=GREEN, header_color=GREEN)

add_footer(s, 7)


# ---- Slide 8: Training picture (curve + 3 numbers) -------------------
s = add_blank(prs)
add_title_band(s, "Training  —  one decision per channel, 7000 channels")

add_figure(s, "13_training_curve_clean.png",
           left=0.5, top=1.2, width=7.4)
add_textbox(s, 0.5, 6.55, 7.4, 0.4,
            "Smoothed running-average reward climbs as ε decays to "
            "0.05; policy stabilises in the second half.",
            size=11, italic=True, color=GREY, align=PP_ALIGN.CENTER)

# Three "stat" cards on the right
add_card(s, 8.2, 1.2, 4.7, 1.7, fill=SOFT_NAVY, border=NAVY)
add_textbox(s, 8.2, 1.35, 4.7, 0.5, "7000",
            size=44, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_textbox(s, 8.2, 2.25, 4.7, 0.5,
            "training episodes per Nt",
            size=14, italic=True, color=DARK, align=PP_ALIGN.CENTER)

add_card(s, 8.2, 3.05, 4.7, 1.7, fill=SOFT_GRN, border=GREEN)
add_textbox(s, 8.2, 3.20, 4.7, 0.5, "≈ 38 s",
            size=44, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_textbox(s, 8.2, 4.10, 4.7, 0.5,
            "training cost on a single CPU core",
            size=14, italic=True, color=DARK, align=PP_ALIGN.CENTER)

add_card(s, 8.2, 4.9, 4.7, 1.7, fill=SOFT_RED, border=RED)
add_textbox(s, 8.2, 5.05, 4.7, 0.5, "Nt ∈ {2, 4, 8}",
            size=28, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_textbox(s, 8.2, 5.95, 4.7, 0.5,
            "three independent models, identical hyperparams",
            size=14, italic=True, color=DARK, align=PP_ALIGN.CENTER)

add_footer(s, 8)


# ---- Slide 9: Headline result (FULL FIGURE) --------------------------
s = add_blank(prs)
add_title_band(s, "Headline result  —  three schemes, same channels",
               "κ = 0.4, Nt = 4, 400 unseen channels per SNR point")

add_figure(s, "04_three_scheme_comparison.png",
           left=2.6, top=1.2, width=8.0)

# Explicit axis-callout strip
add_textbox(s, 0.5, 6.05, 6.2, 0.4,
            "X axis  →  Transmit SNR (dB)",
            size=13, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
add_textbox(s, 6.5, 6.05, 6.3, 0.4,
            "Top Y  →  absolute Rs (bits/s/Hz);   "
            "Bottom Y  →  gain over Fixed",
            size=13, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

add_textbox(s, 0.6, 6.55, 12.3, 0.8,
            "DQN (green) sits at or above the Fixed baseline (grey) "
            "everywhere.  Traditional (red, dashed) drops below.  "
            "Same model, no retuning per SNR.",
            size=13, italic=True, color=DARK, align=PP_ALIGN.CENTER)
add_footer(s, 9)


# ---- Slide 10: κ sweep + outage (TWO FIGURES) ------------------------
s = add_blank(prs)
add_title_band(s, "Robustness  —  the κ sweep and the outage curve")

add_figure(s, "06_kappa_sweep.png", left=0.4, top=1.2, width=6.4)
add_figure(s, "08_secrecy_outage.png", left=7.2, top=1.2, width=5.7)

# Per-panel axis callouts
add_textbox(s, 0.4, 5.85, 6.4, 0.4,
            "X →  CSI quality κ ∈ [0, 1]    "
            "Y →  Rs (bits/s/Hz)",
            size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_textbox(s, 7.2, 5.85, 5.7, 0.4,
            "X →  target rate R₀    "
            "Y →  outage P(Rs < R₀)",
            size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_textbox(s, 0.4, 6.30, 6.4, 0.85,
            "Drag κ from 1 → 0.  Traditional collapses below "
            "Fixed near κ ≈ 0.6.  DQN never does.",
            size=12, italic=True, color=DARK, align=PP_ALIGN.CENTER)
add_textbox(s, 7.2, 6.30, 5.7, 0.85,
            "Outage at the same scenario.  At R₀ = 4 bits/s/Hz, "
            "DQN's outage ≈ 18% vs Traditional's ≈ 26%.",
            size=12, italic=True, color=DARK, align=PP_ALIGN.CENTER)
add_footer(s, 10)


# ---- Slide 11: Live demo cue ----------------------------------------
s = add_blank(prs)
add_title_band(s, "Live demo  —  Streamlit GUI")

add_textbox(s, 1.5, 1.4, 10.3, 0.5,
            "Switching to the browser  —  same model, you drag "
            "the sliders.",
            size=20, italic=True, color=NAVY, align=PP_ALIGN.CENTER)

# Five demo cards in a single row of icons
demos = [
    ("1.", "Live test", "Pick (Nt, SNR, κ).  See three schemes "
                       "compete side-by-side."),
    ("2.", "SNR sweep", "Hold κ, vary SNR.  Watch the curves spread "
                        "or converge."),
    ("3.", "κ sweep",  "Hold SNR, slide κ from 1 → 0.  The headline "
                        "story, live."),
    ("4.", "Outage",   "Empirical CDF.  How often each scheme misses "
                       "the target rate."),
    ("5.", "Geometry", "Conceptual sketch  —  Alice's beam, AN, "
                       "Bob, Eve."),
]
x = 0.5
for n, name, body in demos:
    add_card(s, x, 2.3, 2.45, 4.4, fill=SOFT_NAVY, border=NAVY)
    add_textbox(s, x, 2.45, 2.45, 0.5, n,
                size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s, x, 2.95, 2.45, 0.5, name,
                size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(s, x + 0.15, 3.55, 2.15, 3.0, body,
                size=12, color=DARK, align=PP_ALIGN.CENTER)
    x += 2.55

add_textbox(s, 0.6, 6.95, 12.3, 0.4,
            "▶  streamlit run app/streamlit_app.py     "
            "(default Nt = 8 — where the win is biggest)",
            size=13, italic=True, color=GREY, align=PP_ALIGN.CENTER)
add_footer(s, 11)


# ---- Slide 12: Conclusion (3 takeaways) -----------------------------
s = add_blank(prs)
add_title_band(s, "Takeaways")

add_card(s, 0.6, 1.4, 4.0, 5.5, fill=SOFT_RED, border=RED)
add_textbox(s, 0.6, 1.55, 4.0, 0.5, "1.  Problem",
            size=22, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_textbox(s, 0.8, 2.2, 3.6, 4.5,
            "Plugging noisy intel into the textbook AN-allocation "
            "optimiser produces a confidently-wrong split.\n\n"
            "It actively underperforms the do-nothing baseline once "
            "intel quality drops below ~ 0.6.",
            size=15, color=DARK)

add_card(s, 4.7, 1.4, 4.0, 5.5, fill=SOFT_NAVY, border=NAVY)
add_textbox(s, 4.7, 1.55, 4.0, 0.5, "2.  Fix",
            size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_textbox(s, 4.9, 2.2, 3.6, 4.5,
            "Train a small DQN over random intel-quality scenarios.\n\n"
            "Feed the agent its own intel quality and a beam-alignment "
            "cue, so it knows when to trust the estimate and when to "
            "retreat to a safe default.",
            size=15, color=DARK)

add_card(s, 8.8, 1.4, 4.0, 5.5, fill=SOFT_GRN, border=GREEN)
add_textbox(s, 8.8, 1.55, 4.0, 0.5, "3.  Outcome",
            size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_textbox(s, 9.0, 2.2, 3.6, 4.5,
            "ONE trained policy beats the noisy-CSI optimiser at every "
            "Nt and every κ in the test set.\n\n"
            "Lower outage at the same target rate.\n\n"
            "≈ 38 s training cost.  Single-pass inference.",
            size=15, color=DARK)

add_footer(s, 12)


# ---- Slide 13: Q&A --------------------------------------------------
s = add_blank(prs)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          Inches(0), Inches(0),
                          Inches(13.333), Inches(7.5))
band.line.fill.background()
band.fill.solid()
band.fill.fore_color.rgb = NAVY
add_textbox(s, 1.0, 2.5, 11.3, 1.4,
            "Thank you  —  questions?",
            size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, 1.0, 4.2, 11.3, 0.6,
            "github.com/MuhammadIsmail009/pls-dqn-miso",
            size=20, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_textbox(s, 1.0, 4.95, 11.3, 0.5,
            "Muhammad Ismail   ·   Abubakar Butt   ·   Usman Ali",
            size=15, color=RGBColor(0xE0, 0xE5, 0xEC),
            align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Strip identifying metadata before saving
# ---------------------------------------------------------------------------
cp = prs.core_properties
for attr in ("author", "last_modified_by", "title", "subject", "keywords",
             "category", "comments", "identifier", "language",
             "content_status", "version"):
    try:
        setattr(cp, attr, "")
    except Exception:
        pass

prs.save(OUT_PATH)


# Post-process the saved file: wipe app.xml fields and remove customXml /
# thumbnail (same approach we used for the docx).
import zipfile, shutil, re

tmp_path = OUT_PATH + ".tmp"
DROP_PREFIXES = ("customXml/",)
DROP_FILES    = ("docProps/thumbnail.jpeg",)

with zipfile.ZipFile(OUT_PATH, "r") as zin, \
     zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
    for item in zin.infolist():
        name = item.filename
        if any(name.startswith(p) for p in DROP_PREFIXES):
            continue
        if name in DROP_FILES:
            continue
        data = zin.read(name)

        if name == "[Content_Types].xml":
            xml = data.decode("utf-8", errors="ignore")
            xml = re.sub(
                r"<Override[^/]*PartName=\"/customXml/[^\"]+\"[^/]*/>",
                "", xml)
            xml = re.sub(
                r"<Default[^/]*Extension=\"jpeg\"[^/]*/>", "", xml)
            data = xml.encode("utf-8")

        elif name == "_rels/.rels":
            xml = data.decode("utf-8", errors="ignore")
            xml = re.sub(
                r"<Relationship[^/]*Target=\"docProps/thumbnail\.jpeg\"[^/]*/>",
                "", xml)
            data = xml.encode("utf-8")

        elif name == "ppt/_rels/presentation.xml.rels":
            xml = data.decode("utf-8", errors="ignore")
            xml = re.sub(
                r"<Relationship[^/]*Target=\"\.\./customXml/[^\"]+\"[^/]*/>",
                "", xml)
            data = xml.encode("utf-8")

        elif name == "docProps/app.xml":
            xml = data.decode("utf-8", errors="ignore")
            for tag in ("Application", "Company", "Manager", "AppVersion",
                        "Template", "TotalTime", "PresentationFormat"):
                xml = re.sub(rf"<{tag}>[^<]*</{tag}>",
                             f"<{tag}></{tag}>", xml)
            data = xml.encode("utf-8")

        zout.writestr(item, data)

shutil.move(tmp_path, OUT_PATH)

# Sanity: zero-width / hidden chars
suspicious = ["​", "‌", "‍", "⁠", "﻿"]
hidden = False
with zipfile.ZipFile(OUT_PATH, "r") as z:
    for name in z.namelist():
        if name.endswith(".xml"):
            data = z.read(name).decode("utf-8", errors="ignore")
            for ch in suspicious:
                if ch in data:
                    print(f"[WARN] {name} contains hidden char "
                          f"U+{ord(ch):04X}")
                    hidden = True
if not hidden:
    print("[OK] no zero-width/invisible characters found")

print(f"[SAVE] {OUT_PATH}")
print(f"[SIZE] {os.path.getsize(OUT_PATH) / 1024:.1f} KB")
print(f"[SLIDES] {len(prs.slides)}")
